from __future__ import annotations

import math
import os
from pathlib import Path
from textwrap import wrap

from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(r"E:\code\learning-agent-platform")
OUT_DIR = ROOT / "docs" / "demo-deliverables"
PREVIEW_DIR = OUT_DIR / "ppt-preview"
ASSET_DIR = ROOT / "outputs" / "demo-captures"
TMP_DIR = ROOT / "outputs" / "demo-ppt-build"

OUT_DIR.mkdir(parents=True, exist_ok=True)
PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
TMP_DIR.mkdir(parents=True, exist_ok=True)

W, H = 1920, 1080
NAVY = "#10263F"
NAVY_2 = "#173656"
BLUE = "#2B5D8A"
CYAN = "#3D90B8"
GOLD = "#C27A2A"
SOFT_BG = "#F5F7FB"
PANEL = "#FFFFFF"
BORDER = "#D9E2EF"
MUTED = "#5D6B7B"
TEXT = "#162338"
GREEN = "#3F8A7E"
ORANGE = "#D58A43"
RED = "#C95C5C"

FONT_CANDIDATES = [
    Path(r"C:\Windows\Fonts\NotoSansSC-VF.ttf"),
    Path(r"C:\Windows\Fonts\simhei.ttf"),
    Path(r"C:\Windows\Fonts\msyh.ttc"),
]


def pick_font() -> Path:
    for candidate in FONT_CANDIDATES:
        if candidate.exists():
            return candidate
    raise FileNotFoundError("No Chinese font found.")


FONT_PATH = pick_font()


def f(size: int) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(str(FONT_PATH), size=size)


def rgb(hex_value: str) -> tuple[int, int, int]:
    value = hex_value.lstrip("#")
    return tuple(int(value[i : i + 2], 16) for i in (0, 2, 4))


def canvas(color: str = SOFT_BG) -> Image.Image:
    return Image.new("RGB", (W, H), rgb(color))


def draw_rounded(draw: ImageDraw.ImageDraw, box, fill, outline=None, radius=24, width=2):
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def wrap_text(text: str, font: ImageFont.FreeTypeFont, max_width: int) -> list[str]:
    lines: list[str] = []
    for raw in text.splitlines():
        if not raw.strip():
            lines.append("")
            continue
        parts = raw.split(" ")
        current = ""
        for part in parts:
            trial = part if not current else current + " " + part
            if draw_text_width(trial, font) <= max_width:
                current = trial
            else:
                if current:
                    lines.append(current)
                if draw_text_width(part, font) <= max_width:
                    current = part
                else:
                    # Break long runs like URLs or code tokens.
                    buffer = ""
                    for ch in part:
                        trial2 = buffer + ch
                        if draw_text_width(trial2, font) <= max_width:
                            buffer = trial2
                        else:
                            if buffer:
                                lines.append(buffer)
                            buffer = ch
                    current = buffer
        if current:
            lines.append(current)
    return lines or [""]


def draw_text_width(text: str, font: ImageFont.FreeTypeFont) -> int:
    bbox = font.getbbox(text)
    return bbox[2] - bbox[0]


def draw_multiline(draw: ImageDraw.ImageDraw, xy, text, font, fill, max_width, line_gap=10):
    x, y = xy
    lines = wrap_text(text, font, max_width)
    for line in lines:
        draw.text((x, y), line, font=font, fill=fill)
        y += font.size + line_gap
    return y


def draw_bullets(draw, x, y, width, bullets, font, fill=TEXT, bullet_fill=BLUE, line_gap=12, item_gap=16):
    for bullet in bullets:
        draw.ellipse((x, y + 11, x + 10, y + 21), fill=rgb(bullet_fill))
        y = draw_multiline(draw, (x + 24, y), bullet, font, rgb(fill), width - 28, line_gap=line_gap)
        y += item_gap
    return y


def draw_page_footer(draw: ImageDraw.ImageDraw, page_no: int):
    draw.text((W - 138, H - 46), f"{page_no:02d} / 12", font=f(18), fill=rgb(MUTED))


def add_header(draw, title: str, subtitle: str, page_no: int):
    draw.rectangle((0, 0, W, 88), fill=rgb(NAVY))
    draw.rectangle((0, 88, W, 96), fill=rgb(GOLD))
    draw.text((74, 30), title, font=f(42), fill=rgb(TEXT))
    draw.text((78, 112), subtitle, font=f(24), fill=rgb(MUTED))
    draw_rounded(draw, (1540, 26, 1828, 72), fill=rgb(GOLD), radius=20, width=0)
    draw.text((1596, 35), "MVP 预览版", font=f(22), fill=(255, 255, 255))
    draw_page_footer(draw, page_no)


def add_panel(draw, box, fill=PANEL, outline=BORDER, radius=28, width=2):
    draw_rounded(draw, box, fill=rgb(fill), outline=rgb(outline), radius=radius, width=width)


def paste_screenshot(base: Image.Image, screenshot_name: str, box, padding=0, fit="contain"):
    img = Image.open(ASSET_DIR / screenshot_name).convert("RGB")
    x1, y1, x2, y2 = box
    x1 += padding
    y1 += padding
    x2 -= padding
    y2 -= padding
    target = (x2 - x1, y2 - y1)
    if fit == "cover":
        thumb = ImageOps.fit(img, target, method=Image.Resampling.LANCZOS)
    else:
        thumb = ImageOps.contain(img, target, method=Image.Resampling.LANCZOS)
    px = x1 + (target[0] - thumb.width) // 2
    py = y1 + (target[1] - thumb.height) // 2
    base.paste(thumb, (px, py))


def cover_slide() -> Image.Image:
    im = canvas(NAVY)
    d = ImageDraw.Draw(im)
    d.rectangle((0, 0, W, H), fill=rgb(NAVY))
    d.ellipse((-180, -120, 760, 640), fill=rgb(NAVY_2))
    d.ellipse((1280, 680, 2140, 1540), fill=rgb(NAVY_2))
    add_panel(d, (86, 104, 824, 394), fill="#142C47", outline="#28486C", radius=34)
    d.text((132, 138), "AI 编程学习平台", font=f(54), fill=(255, 255, 255))
    d.text((132, 214), "Learning Agent Platform", font=f(30), fill=(221, 232, 244))
    d.text((132, 292), "组号：____    课程 / 比赛名称：____    日期：____", font=f(30), fill=(245, 245, 245))
    d.text((132, 348), "小组成员：5 人占位，后续由用户填写姓名与学号", font=f(28), fill=(221, 232, 244))
    add_panel(d, (86, 448, 790, 764), fill="#FFFFFF", outline="#2D4F76", radius=30)
    d.text((128, 490), "这是一套 5 分钟以内的精简版答辩稿", font=f(32), fill=rgb(TEXT))
    d.text((128, 552), "• 只保留 12 页核心内容", font=f(26), fill=rgb(TEXT))
    d.text((128, 604), "• 字体统一为微软雅黑或系统等价无衬线字体", font=f(26), fill=rgb(TEXT))
    d.text((128, 656), "• 当前阶段标注为 MVP / preview-only / disabled-by-default", font=f(26), fill=rgb(TEXT))
    draw_rounded(d, (920, 150, 1778, 900), fill=rgb("#F7FAFD"), outline=rgb(BORDER), radius=36, width=2)
    paste_screenshot(im, "01-home.png", (950, 180, 1748, 870), padding=0, fit="cover")
    d.text((954, 830), "演示素材截图", font=f(20), fill=rgb(MUTED))
    return im


def team_slide() -> Image.Image:
    im = canvas()
    d = ImageDraw.Draw(im)
    add_header(d, "小组分工", "5 人占位信息与内部评分留空，方便答辩前直接填写。", 2)
    add_panel(d, (72, 180, 1848, 980))
    cols = [260, 260, 890, 320]
    x0 = 114
    y0 = 238
    row_h = 120
    headers = ["姓名", "学号", "具体负责工作", "组内评分 / 贡献分"]
    rows = [
        ["________", "________", "项目创意设计、需求分析、答辩讲解", "________"],
        ["________", "________", "AI 工具选型、提示词协作、代码生成", "________"],
        ["________", "________", "Web 页面测试、导入预览验收", "________"],
        ["________", "________", "演示视频录制、字幕说明、素材整理", "________"],
        ["________", "________", "PPT 制作、项目总结、答辩材料整理", "________"],
    ]
    # header
    cx = x0
    for i, head in enumerate(headers):
        draw_rounded(d, (cx, y0, cx + cols[i], y0 + 92), fill=rgb(NAVY), outline=rgb(NAVY), radius=18, width=0)
        d.text((cx + 18, y0 + 26), head, font=f(26), fill=(255, 255, 255))
        cx += cols[i]
    # body
    for r, row in enumerate(rows):
        cy = y0 + 110 + r * row_h
        cx = x0
        for i, cell in enumerate(row):
            draw_rounded(d, (cx, cy, cx + cols[i], cy + 96), fill=rgb("#FFFFFF"), outline=rgb(BORDER), radius=16, width=2)
            d.text((cx + 18, cy + 29), cell, font=f(24), fill=rgb(TEXT))
            cx += cols[i]
    d.text((114, 910), "提示：姓名、学号、组内评分/贡献分均保留空位。", font=f(24), fill=rgb(MUTED))
    return im


def background_slide() -> Image.Image:
    im = canvas()
    d = ImageDraw.Draw(im)
    add_header(d, "项目背景", "围绕“读书 - 练习 - 总结 - 行动”的学习链路展开，强调真实需求与约束。", 3)
    add_panel(d, (72, 182, 1186, 980))
    bullets = [
        "编程学习资料分散，阅读、练习、总结经常割裂。",
        "初学者难以把书籍内容转化为下一步练习行动。",
        "AI 工具可以辅助学习路径规划和代码内容理解。",
    ]
    d.text((118, 238), "实际需求", font=f(34), fill=rgb(BLUE))
    draw_bullets(d, 118, 308, 960, bullets, font=f(30), line_gap=10, item_gap=20)
    add_panel(d, (1248, 182, 1848, 980), fill="#F9FBFD")
    d.text((1288, 240), "目标用户", font=f(34), fill=rgb(BLUE))
    draw_bullets(
        d,
        1288,
        314,
        470,
        [
            "编程初学者",
            "课程答辩/展示场景",
            "需要安全预览的内容导入场景",
        ],
        font=f(28),
        item_gap=32,
    )
    add_panel(d, (1248, 640, 1848, 980), fill="#EEF5FB")
    d.text((1288, 692), "项目目标", font=f(34), fill=rgb(BLUE))
    draw_bullets(
        d,
        1288,
        758,
        470,
        [
            "先做可运行闭环，再做复杂智能化。",
            "先做 preview-only，再做真实写入。",
        ],
        font=f(28),
        item_gap=24,
    )
    return im


def creativity_slide() -> Image.Image:
    im = canvas()
    d = ImageDraw.Draw(im)
    add_header(d, "核心创意与课程契合度", "把创意、功能价值和课程要求合并到同一条叙事线上。", 4)
    cards = [
        (84, 230, 550, 650, "创意", "编程学习网站 + 阅读器 + 文本导入预览 + Desktop 预览 + AI Agent 规划"),
        (684, 230, 550, 650, "解决问题", "资料阅读、代码块定位、导入整理、学习行动提示"),
        (1284, 230, 552, 650, "课程契合", "AI 工具辅助创意、开发、测试、文档、演示制作"),
    ]
    for x, y, w, h, title, body in cards:
        add_panel(d, (x, y, x + w, y + h))
        draw_rounded(d, (x + 32, y + 34, x + 230, y + 88), fill=rgb(GOLD), outline=rgb(GOLD), radius=18, width=0)
        d.text((x + 58, y + 46), title, font=f(28), fill=(255, 255, 255))
        d.text((x + 34, y + 132), body, font=f(30), fill=rgb(TEXT))
    add_panel(d, (84, 918, 1840, 1016), fill="#F3F8FC")
    d.text((118, 948), "强调：必要性在于解决“读不下去、练不起来、总结不成链”的问题；创新性在于把 AI 与学习路径拆成可控预览层。", font=f(26), fill=rgb(MUTED))
    return im


def tools_slide() -> Image.Image:
    im = canvas()
    d = ImageDraw.Draw(im)
    add_header(d, "AI 工具与模型选择", "列出项目中使用或辅助使用的工具，强调“协作工具”而不是生产能力承诺。", 5)
    cards = [
        (86, 226, 810, 300, "ChatGPT / GPT-5.5", "需求拆解、方案设计、答辩材料规划"),
        (1022, 226, 810, 300, "Codex / GPT-5.4 mini", "代码生成、测试修复、功能迭代"),
        (86, 586, 810, 300, "Claude Code", "辅助提交、代码审查或局部执行"),
        (1022, 586, 810, 300, "DeepSeek 本地文档 Agent", "轮次交接、阶段总结、上下文压缩"),
    ]
    for x, y, w, h, title, body in cards:
        add_panel(d, (x, y, x + w, y + h))
        draw_rounded(d, (x + 28, y + 30, x + 330, y + 84), fill=rgb(BLUE), outline=rgb(BLUE), radius=18, width=0)
        d.text((x + 54, y + 42), title, font=f(28), fill=(255, 255, 255))
        d.text((x + 34, y + 126), body, font=f(30), fill=rgb(TEXT))
    add_panel(d, (86, 936, 1844, 1008), fill="#F7FBFF")
    d.text((120, 958), "说明：以上工具用于创意、开发、测试与文档协作，不等于系统生产能力。", font=f(24), fill=rgb(MUTED))
    return im


def workflow_slide() -> Image.Image:
    im = canvas()
    d = ImageDraw.Draw(im)
    add_header(d, "AI 工具运用流程", "用流程图表达“提需求 - 生成提示 - 修改代码 - handoff - 人工验收”的闭环。", 6)
    steps = [
        "1\n提出功能目标",
        "2\nChatGPT 生成 Codex 提示词",
        "3\nCodex 修改代码并运行测试",
        "4\nDeepSeek 生成 handoff 文档",
        "5\n人工验收与下一轮迭代",
    ]
    xs = [92, 430, 768, 1106, 1444]
    colors = [BLUE, GOLD, CYAN, GREEN, ORANGE]
    for i, (x, label) in enumerate(zip(xs, steps)):
        add_panel(d, (x, 380, x + 280, 560), fill="#FFFFFF")
        draw_rounded(d, (x + 98, 324, x + 182, 408), fill=rgb(colors[i]), outline=rgb(colors[i]), radius=42, width=0)
        d.text((x + 126, 344), str(i + 1), font=f(28), fill=(255, 255, 255))
        d.text((x + 48, 442), label, font=f(28), fill=rgb(TEXT), align="center")
        if i < 4:
            d.line((x + 280, 470, x + 336, 470), fill=rgb(BORDER), width=6)
            d.polygon([(x + 336, 470), (x + 318, 458), (x + 318, 482)], fill=rgb(BORDER))
    add_panel(d, (132, 664, 1788, 932), fill="#F7FAFD")
    d.text((174, 708), "流程要点", font=f(34), fill=rgb(BLUE))
    draw_bullets(
        d,
        178,
        772,
        1480,
        [
            "少文字、少范围，先把一个明确功能做成闭环。",
            "每轮都有人审查，避免 AI 自动扩大任务边界。",
            "handoff 文档用于压缩上下文、保存阶段结论。",
        ],
        font=f(28),
        line_gap=8,
        item_gap=12,
    )
    return im


def ai_role_slide() -> Image.Image:
    im = canvas()
    d = ImageDraw.Draw(im)
    add_header(d, "AI 技术在项目中的核心作用", "AI 主要承担协作、生成、验证和总结，不替代人工验收。", 7)
    bullets = [
        "生成 Next.js / TypeScript / Electron 代码",
        "设计 preview-only 安全边界",
        "生成测试用例并修复类型错误",
        "生成演示视频脚本和 PPT 初稿",
        "辅助总结项目进度与风险",
    ]
    add_panel(d, (86, 220, 1280, 950))
    draw_bullets(d, 124, 286, 1088, bullets, font=f(32), item_gap=26)
    add_panel(d, (1340, 220, 1840, 950), fill="#F7FBFD")
    d.text((1380, 284), "原则", font=f(34), fill=rgb(BLUE))
    draw_bullets(
        d,
        1380,
        354,
        380,
        [
            "AI 是协作工具",
            "人工验收必须存在",
            "边界控制优先",
            "安全审查优先",
        ],
        font=f(28),
        item_gap=28,
    )
    add_panel(d, (86, 980, 1840, 1030), fill="#EEF4FA")
    d.text((122, 990), "结论：真实项目仍需要人工验收、边界控制和安全审查，AI 只是加速工具。", font=f(24), fill=rgb(MUTED))
    return im


def problem_solution_slide() -> Image.Image:
    im = canvas()
    d = ImageDraw.Draw(im)
    add_header(d, "AI 工具使用中的问题与解决", "用“问题—解决”的方式快速说明本项目在协作中的应对策略。", 8)
    add_panel(d, (88, 214, 1838, 958))
    x1, y1, w1, h = 124, 286, 730, 100
    x2 = 992
    header_h = 80
    d.rounded_rectangle((x1, 240, x1 + 760, 320), radius=18, fill=rgb(NAVY), outline=rgb(NAVY))
    d.text((x1 + 28, 259), "问题", font=f(28), fill=(255, 255, 255))
    d.rounded_rectangle((x2, 240, x2 + 760, 320), radius=18, fill=rgb(NAVY), outline=rgb(NAVY))
    d.text((x2 + 28, 259), "解决", font=f(28), fill=(255, 255, 255))
    rows = [
        ("上下文过长", "DeepSeek handoff / 阶段压缩"),
        ("AI 容易扩大范围", "每轮限制允许/禁止目录"),
        ("安全风险", "preview-only、disabled-by-default、不写 DB"),
        ("PPT 文本溢出", "重新排版、减少文字、统一字号"),
        ("Git/push 网络问题", "本地确认 commit，网络恢复后单独 push"),
    ]
    y = 340
    for problem, solution in rows:
        for x, text in [(x1, problem), (x2, solution)]:
            add_panel(d, (x, y, x + 760, y + 88), fill="#FFFFFF")
            d.text((x + 26, y + 24), text, font=f(28), fill=rgb(TEXT))
        y += 104
    return im


def reader_slide() -> Image.Image:
    im = canvas()
    d = ImageDraw.Draw(im)
    add_header(d, "成果展示一：Web Reader", "Books / Reader 阅读链路 + 代码块识别 / 目录能力。", 9)
    add_panel(d, (80, 220, 1174, 958))
    paste_screenshot(im, "04-reader-top.png", (106, 250, 1148, 928), fit="cover")
    add_panel(d, (1248, 220, 1848, 958), fill="#F8FBFD")
    d.text((1286, 274), "说明", font=f(34), fill=rgb(BLUE))
    draw_bullets(
        d,
        1288,
        350,
        470,
        [
            "Books / Reader 阅读链路已经连通。",
            "代码块目录、语言筛选、键盘高亮逻辑已实现。",
            "当前开发数据未命中稳定样例，因此用导入页的安全脱敏片段补位说明。",
        ],
        font=f(28),
        item_gap=26,
    )
    add_panel(d, (1248, 670, 1848, 958), fill="#EEF5FB")
    d.text((1286, 708), "简单示意", font=f(30), fill=rgb(BLUE))
    d.text((1288, 772), "代码块目录  →  语言筛选  →  键盘定位高亮", font=f(26), fill=rgb(TEXT))
    d.text((1288, 826), "提升代码类教材的阅读体验。", font=f(26), fill=rgb(TEXT))
    return im


def import_slide() -> Image.Image:
    im = canvas()
    d = ImageDraw.Draw(im)
    add_header(d, "成果展示二：纯文本导入预览", "粘贴文本、自动章节切分、草案编辑、保存按钮 disabled。", 10)
    add_panel(d, (80, 220, 1220, 958))
    paste_screenshot(im, "08-import-chapter-edit.png", (112, 254, 1188, 926), fit="cover")
    d.rounded_rectangle((104, 846, 212, 926), radius=18, fill=rgb("#F8FBFD"), outline=rgb("#F8FBFD"))
    add_panel(d, (1284, 220, 1848, 958), fill="#F8FBFD")
    draw_bullets(
        d,
        1318,
        292,
        448,
        [
            "支持粘贴纯文本与 Markdown 标题。",
            "自动章节切分后先生成确认草案。",
            "章节可重命名、排除 / 恢复、undo / redo。",
            "保存按钮 disabled，当前未写 DB。",
        ],
        font=f(28),
        item_gap=22,
    )
    add_panel(d, (1284, 700, 1848, 958), fill="#EEF5FB")
    d.text((1318, 744), "阶段说明", font=f(30), fill=rgb(BLUE))
    d.text((1318, 804), "当前是 preview-only，", font=f(26), fill=rgb(TEXT))
    d.text((1318, 854), "为后续真实导入做准备。", font=f(26), fill=rgb(TEXT))
    return im


def desktop_slide() -> Image.Image:
    im = canvas()
    d = ImageDraw.Draw(im)
    add_header(d, "成果展示三：Desktop 本地预览", "阅读进度、书签和今日学习行动卡片都只做本地展示。", 11)
    add_panel(d, (80, 220, 1220, 958))
    paste_screenshot(im, "11-desktop-home.png", (112, 254, 1188, 926), fit="cover")
    add_panel(d, (1284, 220, 1848, 958), fill="#F8FBFD")
    draw_bullets(
        d,
        1318,
        292,
        448,
        [
            "本地阅读进度",
            "书签预览",
            "今日学习行动本地预览",
        ],
        font=f(30),
        item_gap=38,
    )
    add_panel(d, (1284, 700, 1848, 958), fill="#EEF5FB")
    d.text((1318, 744), "说明", font=f(30), fill=rgb(BLUE))
    d.text((1318, 804), "Desktop 端目前是 local-only / preview-only，", font=f(24), fill=rgb(TEXT))
    d.text((1318, 854), "只读展示，不连接真实后端。", font=f(24), fill=rgb(TEXT))
    return im


def summary_slide() -> Image.Image:
    im = canvas()
    d = ImageDraw.Draw(im)
    add_header(d, "总结与反思", "把长期复杂项目拆成可展示、可验证、可继续推进的阶段成果。", 12)
    blocks = [
        (86, 224, 820, 300, "收获", ["学习了 AI 协作开发、前后端架构、测试和演示材料制作。"], GREEN),
        (1014, 224, 820, 300, "挑战", ["项目范围大、上下文管理难、安全边界复杂、UI 美化不足。"], ORANGE),
        (86, 572, 820, 300, "解决", ["小轮迭代、handoff 文档、preview-only 机制、测试验证。"], BLUE),
        (1014, 572, 820, 300, "未来改进", ["真实登录、DB 保存、LLM/RAG、PDF/EPUB/URL 导入、Desktop 深度集成、UI 美化。"], CYAN),
    ]
    for x, y, w, h, title, lines, color in blocks:
        add_panel(d, (x, y, x + w, y + h))
        draw_rounded(d, (x + 32, y + 28, x + 190, y + 82), fill=rgb(color), outline=rgb(color), radius=18, width=0)
        d.text((x + 64, y + 40), title, font=f(28), fill=(255, 255, 255))
        draw_bullets(d, x + 34, y + 118, w - 68, lines, font=f(28), item_gap=20)
    add_panel(d, (86, 918, 1834, 1016), fill="#F3F8FC")
    d.text((122, 948), "当前状态已明确标注为 MVP / preview-only / disabled-by-default，不是生产上线。", font=f(24), fill=rgb(MUTED))
    return im


SLIDES = [
    cover_slide,
    team_slide,
    background_slide,
    creativity_slide,
    tools_slide,
    workflow_slide,
    ai_role_slide,
    problem_solution_slide,
    reader_slide,
    import_slide,
    desktop_slide,
    summary_slide,
]


def build_preview_images() -> list[Path]:
    paths: list[Path] = []
    for idx, fn in enumerate(SLIDES, start=1):
        img = fn()
        preview_path = PREVIEW_DIR / f"slide-{idx:02d}.png"
        img.save(preview_path)
        paths.append(preview_path)
    return paths


def build_pptx(slide_images: list[Path], output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for img_path in slide_images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(img_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(str(output_path))


def main() -> None:
    slide_images = build_preview_images()
    build_pptx(slide_images, OUT_DIR / "defense-presentation-revised.pptx")


if __name__ == "__main__":
    main()
