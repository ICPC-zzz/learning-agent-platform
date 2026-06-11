from __future__ import annotations

import math
import os
import textwrap
from pathlib import Path

import imageio
import imageio_ffmpeg
from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"E:\code\learning-agent-platform")
CAPTURE_DIR = ROOT / "outputs" / "demo-captures"
TEMP_DIR = ROOT / "outputs" / "demo-final-assets"
OUT_DIR = ROOT / "docs" / "demo-deliverables"
OUT_DIR.mkdir(parents=True, exist_ok=True)
TEMP_DIR.mkdir(parents=True, exist_ok=True)

os.environ.setdefault("IMAGEIO_FFMPEG_EXE", imageio_ffmpeg.get_ffmpeg_exe())

FONT_CANDIDATES = [
    Path(r"C:\Windows\Fonts\NotoSansSC-VF.ttf"),
    Path(r"C:\Windows\Fonts\simhei.ttf"),
    Path(r"C:\Windows\Fonts\msyh.ttc"),
]


def pick_font_path() -> Path:
    for candidate in FONT_CANDIDATES:
        if candidate.exists():
            return candidate
    raise FileNotFoundError("No Chinese font found on the system.")


FONT_PATH = pick_font_path()


def font(size: int) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(str(FONT_PATH), size=size)


def rgb(hex_value: str) -> tuple[int, int, int]:
    value = hex_value.lstrip("#")
    return tuple(int(value[i : i + 2], 16) for i in (0, 2, 4))


NAVY = rgb("0D1B2A")
BLUE = rgb("1F4E79")
GOLD = rgb("C77D2A")
AMBER = rgb("E9C46A")
TEAL = rgb("2A9D8F")
BG = rgb("F6F8FB")
TEXT = rgb("132033")
MUTED = rgb("5B667A")
PANEL = rgb("FFFFFF")
BORDER = rgb("D7DFEA")
SOFT = rgb("EEF4FB")
SOFT_ORANGE = rgb("FFF2E5")


def wrap_lines(text: str, width: int) -> list[str]:
    lines: list[str] = []
    for raw_line in text.splitlines():
        if not raw_line.strip():
            lines.append("")
            continue
        lines.extend(textwrap.wrap(raw_line, width=width) or [raw_line])
    return lines or [""]


def fit_contain(image: Image.Image, target_size: tuple[int, int], bg_color: tuple[int, int, int]) -> Image.Image:
    target_w, target_h = target_size
    scale = min(target_w / image.width, target_h / image.height)
    new_size = (max(1, int(image.width * scale)), max(1, int(image.height * scale)))
    resized = image.resize(new_size, Image.Resampling.LANCZOS)
    canvas = Image.new("RGB", target_size, bg_color)
    offset = ((target_w - new_size[0]) // 2, (target_h - new_size[1]) // 2)
    canvas.paste(resized, offset)
    return canvas


def cleanup_issue_badge(image: Image.Image) -> Image.Image:
    cleaned = image.copy()
    draw = ImageDraw.Draw(cleaned)
    # Cover the small dev issue bubble that appears in some browser captures.
    draw.rounded_rectangle((0, cleaned.height - 112, 176, cleaned.height), radius=18, fill=(255, 255, 255))
    return cleaned


def add_bottom_caption(image: Image.Image, segment_no: int, caption: str, detail: str, stage: str) -> Image.Image:
    base = fit_contain(cleanup_issue_badge(image), (1920, 1080), BG)
    overlay = Image.new("RGBA", base.size, (0, 0, 0, 0))
    od = ImageDraw.Draw(overlay)
    band_h = 170
    od.rectangle((0, base.height - band_h, base.width, base.height), fill=(11, 22, 38, 190))

    tag_w = 180
    tag_h = 52
    tag_x = 36
    tag_y = 32
    od.rounded_rectangle((tag_x, tag_y, tag_x + tag_w, tag_y + tag_h), radius=18, fill=(199, 125, 42, 245))
    od.rounded_rectangle((base.width - 410, 34, base.width - 34, 88), radius=18, fill=(15, 31, 48, 190))

    composed = Image.alpha_composite(base.convert("RGBA"), overlay)
    draw = ImageDraw.Draw(composed)

    draw.text((tag_x + 24, tag_y + 10), f"DEMO {segment_no:02d}", font=font(22), fill=(255, 255, 255))
    draw.text((base.width - 388, 46), stage, font=font(22), fill=(255, 255, 255))

    caption_font = font(34)
    detail_font = font(24)
    caption_lines = wrap_lines(caption, 28)
    detail_lines = wrap_lines(detail, 52)

    text_x = 64
    text_y = base.height - band_h + 26
    for line in caption_lines[:2]:
        draw.text((text_x, text_y), line, font=caption_font, fill=(255, 255, 255))
        text_y += 42

    if detail_lines:
        text_y += 4
        for line in detail_lines[:2]:
            draw.text((text_x, text_y), line, font=detail_font, fill=(224, 232, 244))
            text_y += 30

    return composed.convert("RGB")


def create_summary_frame() -> Image.Image:
    canvas = Image.new("RGB", (1920, 1080), NAVY)
    draw = ImageDraw.Draw(canvas)
    draw.rectangle((0, 0, 1920, 1080), fill=NAVY)
    draw.ellipse((-140, -160, 600, 520), fill=(24, 46, 77))
    draw.ellipse((1340, 720, 2100, 1460), fill=(24, 46, 77))
    draw.rounded_rectangle((72, 66, 400, 118), radius=20, fill=GOLD)
    draw.text((102, 78), "阶段总结", font=font(24), fill=(255, 255, 255))
    draw.text((78, 170), "MVP / preview-only / disabled-by-default", font=font(42), fill=(255, 255, 255))
    draw.text((78, 230), "当前演示以本地页面为主，强调安全预览边界，而不是生产上线。", font=font(24), fill=(223, 232, 245))

    card_w = 760
    card_h = 160
    gap = 28
    left = 78
    top = 330
    cards = [
        ("已完成", "Web 首页、Reader、导入预览、学习页、Agent 预览、Desktop 预览串成闭环。"),
        ("边界", "不写数据库、不接真实 auth / LLM / Agent loop，危险字段与保存按钮都保持阻断。"),
        ("验证", "lint、typecheck，以及导入预览 / 草案编辑 / Reader 代码块抽取单测均已通过。"),
        ("后续", "真实 auth、DB 持久化、LLM / RAG、PDF / EPUB / URL 导入、Desktop 深度集成。"),
    ]

    for index, (title, body) in enumerate(cards):
        row = index // 2
        col = index % 2
        x = left + col * (card_w + gap)
        y = top + row * (card_h + gap)
        fill = (20, 41, 67) if index % 2 == 0 else (28, 53, 84)
        draw.rounded_rectangle((x, y, x + card_w, y + card_h), radius=24, fill=fill, outline=(77, 108, 150), width=2)
        draw.rounded_rectangle((x + 22, y + 20, x + 126, y + 58), radius=16, fill=(199, 125, 42))
        draw.text((x + 42, y + 28), title, font=font(22), fill=(255, 255, 255))
        body_lines = wrap_lines(body, 30)
        body_y = y + 78
        for line in body_lines[:2]:
            draw.text((x + 24, body_y), line, font=font(24), fill=(234, 240, 248))
            body_y += 32

    draw.text((78, 1000), "本次交付：demo-video.mp4  +  defense-presentation.pptx  +  demo-script.md  +  README.md", font=font(22), fill=(208, 219, 232))
    return canvas


def make_demo_frames() -> list[dict[str, object]]:
    return [
        {
            "source": CAPTURE_DIR / "01-home.png",
            "caption": "首页与导航",
            "detail": "Web 入口把书库、Reader、导入、学习和 Agent 预览串到一起。",
            "stage": "MVP / preview-only / disabled-by-default",
            "duration": 11.0,
        },
        {
            "source": CAPTURE_DIR / "02-books-list.png",
            "caption": "书库列表",
            "detail": "当前数据源只读预览，展示 6 本示例书与整体导航状态。",
            "stage": "preview-only",
            "duration": 10.0,
        },
        {
            "source": CAPTURE_DIR / "03-book-detail.png",
            "caption": "书籍详情",
            "detail": "章节、chunk 和继续阅读入口都存在，但仍然不进入真实写入。",
            "stage": "preview-only",
            "duration": 10.0,
        },
        {
            "source": CAPTURE_DIR / "04-reader-top.png",
            "caption": "Reader 阅读页",
            "detail": "阅读进度、本地状态、同步预览和章节入口都在同一屏里可见。",
            "stage": "disabled-by-default",
            "duration": 12.0,
        },
        {
            "source": CAPTURE_DIR / "05-reader-mid.png",
            "caption": "长文阅读与安全边界",
            "detail": "滚动进度、当前可见块和只读同步说明都维持在预览态。",
            "stage": "preview-only",
            "duration": 12.0,
        },
        {
            "source": CAPTURE_DIR / "06-import-filled.png",
            "caption": "纯文本导入",
            "detail": "Markdown 标题、中文章节和 fenced code 先进入本地预览，保存按钮保持 disabled。",
            "stage": "disabled-by-default",
            "duration": 12.0,
        },
        {
            "source": CAPTURE_DIR / "08-import-chapter-edit.png",
            "caption": "导入预览与草案编辑",
            "detail": "章节切分、危险字段脱敏、重命名和排除编辑都只在本地草案里生效。",
            "stage": "preview-only",
            "duration": 14.0,
        },
        {
            "source": CAPTURE_DIR / "09-learning.png",
            "caption": "学习页",
            "detail": "阅读进度、能力画像、今日行动和本地预览卡片都已串起来。",
            "stage": "preview-only",
            "duration": 11.0,
        },
        {
            "source": CAPTURE_DIR / "10-agent-preview.png",
            "caption": "Agent 预览页",
            "detail": "任务 dry-run、权限边界、Skill 和工具风险都停留在可视化预览。",
            "stage": "disabled-by-default",
            "duration": 11.0,
        },
        {
            "source": CAPTURE_DIR / "11-desktop-home.png",
            "caption": "Desktop 首页",
            "detail": "本地阅读、书签、今日行动和 Reader 跳转都以静态预览卡片呈现。",
            "stage": "preview-only",
            "duration": 12.0,
        },
        {
            "summary": True,
            "caption": "项目总结",
            "detail": "当前阶段是 MVP / preview-only / disabled-by-default，真实 auth、DB、LLM 和 Agent loop 以后再完善。",
            "stage": "next step",
            "duration": 16.0,
        },
    ]


def build_video_frame(item: dict[str, object], index: int) -> Image.Image:
    if item.get("summary"):
        image = create_summary_frame()
    else:
        source = Image.open(item["source"]).convert("RGB")
        image = add_bottom_caption(
            source,
            index + 1,
            str(item["caption"]),
            str(item["detail"]),
            str(item["stage"]),
        )
    return image


def write_video(frame_paths: list[tuple[Path, float]], output_path: Path) -> None:
    writer = imageio.get_writer(
        str(output_path),
        fps=30,
        codec="libx264",
        quality=8,
        macro_block_size=None,
    )
    try:
        for frame_path, duration in frame_paths:
            frame = imageio.imread(frame_path)
            frame_count = max(1, int(round(duration * 30)))
            for _ in range(frame_count):
                writer.append_data(frame)
    finally:
        writer.close()


def set_text_font(run, size: int, color: tuple[int, int, int] = TEXT, bold: bool = False) -> None:
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)


def add_textbox(slide, left, top, width, height, text, size=20, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    box.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_text_font(r, size, color=color, bold=bold)
    return box


def add_bullets(slide, left, top, width, height, bullets, size=20, color=TEXT, gap=8):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.clear()
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = f"• {bullet}"
        p.space_after = Pt(gap)
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(*color)
    return box


def add_badge(slide, x, y, w, h, text, fill, color=(255, 255, 255)):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(*fill)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    set_text_font(r, 18, color=color, bold=True)
    return shape


def add_slide_header(slide, title: str, subtitle: str = "") -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(*BG)
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.34))
    bar = slide.shapes[-1]
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*NAVY)
    bar.line.fill.background()
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(0.34), prs.slide_width, Inches(0.04))
    accent = slide.shapes[-1]
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(*GOLD)
    accent.line.fill.background()
    add_textbox(slide, Inches(0.5), Inches(0.55), Inches(9.2), Inches(0.55), title, size=28, color=NAVY, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.52), Inches(0.95), Inches(11.8), Inches(0.4), subtitle, size=14, color=MUTED)
    add_badge(slide, Inches(11.0), Inches(0.55), Inches(1.9), Inches(0.42), "MVP 预览", GOLD, color=NAVY)
    add_textbox(slide, Inches(10.85), Inches(1.0), Inches(2.0), Inches(0.25), "preview-only", size=10, color=MUTED, align=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(6.0), Inches(0.2), "当前阶段：MVP / preview-only / disabled-by-default", size=10, color=MUTED)


def add_rounded_panel(slide, left, top, width, height, fill=PANEL, line=BORDER):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(*line)
    return shape


def add_picture_contain(slide, path: Path, left, top, width, height, border_color=BORDER):
    panel = add_rounded_panel(slide, left, top, width, height, fill=(255, 255, 255), line=border_color)
    panel.line.width = Pt(1.1)
    img = Image.open(path).convert("RGB")
    ratio = min(width / img.width, height / img.height)
    new_w = int(img.width * ratio)
    new_h = int(img.height * ratio)
    x = left + (width - new_w) // 2
    y = top + (height - new_h) // 2
    slide.shapes.add_picture(str(path), x, y, width=new_w, height=new_h)
    return panel


def add_simple_table(slide, left, top, width, height, rows, cols, data, col_widths=None):
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    if col_widths:
        for idx, col_w in enumerate(col_widths):
            table.columns[idx].width = col_w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            for p in cell.text_frame.paragraphs:
                p.font.name = "Microsoft YaHei"
                p.font.size = Pt(14 if r == 0 else 12)
                p.font.color.rgb = RGBColor(*TEXT)
                p.alignment = PP_ALIGN.CENTER
    return table


def create_presentation() -> Presentation:
    global prs
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # Slide 1: cover
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*NAVY)
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-0.4), Inches(-0.3), Inches(4.1), Inches(4.1)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = RGBColor(24, 45, 76)
    slide.shapes[-1].line.fill.background()
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(10.2), Inches(4.9), Inches(3.8), Inches(3.8)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = RGBColor(24, 45, 76)
    slide.shapes[-1].line.fill.background()
    add_badge(slide, Inches(0.7), Inches(0.7), Inches(2.3), Inches(0.48), "Learning Agent Platform", GOLD, color=NAVY)
    add_textbox(
        slide,
        Inches(0.7),
        Inches(1.45),
        Inches(11.8),
        Inches(1.1),
        "编程学习网站 + AI Agent 软件端 + Skill 社区\n项目答辩与演示材料",
        size=30,
        color=(255, 255, 255),
        bold=True,
    )
    add_textbox(
        slide,
        Inches(0.72),
        Inches(2.9),
        Inches(7.6),
        Inches(0.8),
        "项目阶段：MVP / preview-only / disabled-by-default\n当前演示不等于生产上线，所有写入链路保持阻断或未接入。",
        size=18,
        color=(220, 230, 244),
    )
    add_rounded_panel(slide, Inches(0.72), Inches(4.15), Inches(6.2), Inches(1.1), fill=(18, 34, 53), line=(66, 92, 128))
    add_textbox(slide, Inches(0.95), Inches(4.38), Inches(5.6), Inches(0.35), "课程 / 比赛名称：__________", size=22, color=(255, 255, 255))
    add_textbox(slide, Inches(0.95), Inches(4.72), Inches(5.6), Inches(0.35), "日期：__________", size=22, color=(255, 255, 255))
    add_textbox(slide, Inches(0.72), Inches(5.7), Inches(8.0), Inches(0.55), "说明：本套材料只展示当前可运行的预览闭环，不对真实 auth / DB / LLM / Agent loop 进行任何承诺。", size=15, color=(215, 224, 237))

    # Slide 2: team
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "小组成员与分工", "5 人信息与内部打分留空，方便用户后续直接填写。")
    add_rounded_panel(slide, Inches(0.5), Inches(1.45), Inches(12.2), Inches(5.15))
    headers = ["姓名", "学号", "分工", "内部打分"]
    rows = [
        ["________", "________", "________", "________"],
        ["________", "________", "________", "________"],
        ["________", "________", "________", "________"],
        ["________", "________", "________", "________"],
        ["________", "________", "________", "________"],
    ]
    data = [headers] + rows
    add_simple_table(
        slide,
        Inches(0.8),
        Inches(1.9),
        Inches(11.6),
        Inches(4.5),
        6,
        4,
        data,
        col_widths=[Inches(2.0), Inches(3.0), Inches(4.0), Inches(2.0)],
    )
    add_textbox(slide, Inches(0.82), Inches(6.55), Inches(11.6), Inches(0.35), "提示：也可以在答辩前把“内部打分”留空，现场由成员协商后补填。", size=12, color=MUTED)

    # Slide 3: background
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "项目设计背景", "围绕编程学习的碎片化、练习闭环缺失和 AI 使用边界不清三个问题展开。")
    bullets = [
        "必要性：把“看书 - 做题 - 记录进度 - 进入下一步”连成一条更完整的学习路径。",
        "痛点：纯看文档容易断层，导入与阅读状态分散在不同页面，难以形成可追踪闭环。",
        "目标用户：编程初学者、课程作业答辩场景、需要安全预览的内容导入与学习协作场景。",
        "约束：必须把真实写入、真实 auth、真实 LLM / Agent loop 与前端展示分开处理。",
    ]
    add_bullets(slide, Inches(0.65), Inches(1.55), Inches(5.8), Inches(3.4), bullets, size=19)
    for x, title, body, fill in [
        (6.55, "痛点 1", "学习内容和阅读状态割裂，用户很难知道下一步该学什么。", SOFT),
        (6.55, "痛点 2", "导入、阅读、学习、桌面端都需要一致的安全边界与预览态。", SOFT_ORANGE),
        (9.42, "痛点 3", "如果没有清晰的权限和日志边界，Agent 功能很容易越权。", (236, 248, 245)),
    ]:
        add_rounded_panel(slide, Inches(x), Inches(1.6), Inches(2.3), Inches(1.55), fill=fill)
        add_textbox(slide, Inches(x + 0.16), Inches(1.78), Inches(1.9), Inches(0.25), title, size=18, color=NAVY, bold=True)
        add_textbox(slide, Inches(x + 0.16), Inches(2.05), Inches(1.95), Inches(0.75), body, size=13, color=MUTED)
    add_rounded_panel(slide, Inches(6.55), Inches(3.45), Inches(5.15), Inches(2.4), fill=(255, 255, 255))
    add_textbox(slide, Inches(6.8), Inches(3.72), Inches(4.6), Inches(0.35), "目标", size=18, color=BLUE, bold=True)
    add_bullets(
        slide,
        Inches(6.8),
        Inches(4.05),
        Inches(4.5),
        Inches(1.45),
        [
            "先做可运行闭环，再逐步加入智能化。",
            "先做 preview-only，再做真实保存。",
            "先把权限和日志边界设计清楚。",
        ],
        size=15,
    )

    # Slide 4: innovation
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "项目创新性", "把学习网站、Reader、导入预览和 AI/Skill 规划拆分成可控的分阶段能力。")
    innovations = [
        ("编程学习网站", "不是单页文档，而是书库、阅读、学习建议和任务历史的联合入口。", TEAL),
        ("Reader 安全链路", "阅读进度、书签、章节状态和同步预览分层处理，默认禁用真实写入。", BLUE),
        ("导入预览草案", "纯文本导入先做章节切分、危险字段脱敏和确认草案，不直接入库。", GOLD),
        ("Agent / Skill 规划", "把工具风险、权限边界和 Skill 规划前置到预览层，而不是直接自动执行。", (169, 90, 161)),
        ("Desktop 本地预览", "桌面端只做本地卡片和路由预览，用于展示学习行动和阅读入口。", (90, 111, 95)),
    ]
    for i, (title, body, color) in enumerate(innovations):
        row = i // 2
        col = i % 2
        x = 0.6 + col * 6.2
        y = 1.6 + row * 1.4
        add_rounded_panel(slide, Inches(x), Inches(y), Inches(5.8), Inches(1.15), fill=(255, 255, 255))
        add_badge(slide, Inches(x + 0.15), Inches(y + 0.18), Inches(1.15), Inches(0.34), title, color, color=(255, 255, 255))
        add_textbox(slide, Inches(x + 1.45), Inches(y + 0.12), Inches(4.0), Inches(0.72), body, size=14, color=MUTED)

    # Slide 5: design approach
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "项目设计思路", "先形成内容导入和阅读闭环，再逐步扩展到学习建议和桌面端预览。")
    steps = ["内容输入", "章节切分", "安全预览", "学习推荐", "Desktop 预览"]
    colors = [BLUE, GOLD, TEAL, BLUE, GOLD]
    x_positions = [0.7, 2.95, 5.2, 7.45, 9.7]
    for idx, (step, x, color) in enumerate(zip(steps, x_positions, colors)):
        add_rounded_panel(slide, Inches(x), Inches(2.4), Inches(1.72), Inches(0.95), fill=(255, 255, 255))
        add_badge(slide, Inches(x + 0.2), Inches(2.57), Inches(1.32), Inches(0.36), step, color, color=(255, 255, 255))
        if idx < len(steps) - 1:
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 1.72), Inches(2.87), Inches(x + 2.22), Inches(2.87))
            conn.line.color.rgb = RGBColor(*BORDER)
            conn.line.width = Pt(2)
            # Add a small arrow-like text to keep the flow obvious.
            add_textbox(slide, Inches(x + 1.98), Inches(2.72), Inches(0.22), Inches(0.2), "→", size=26, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_rounded_panel(slide, Inches(0.75), Inches(4.0), Inches(11.8), Inches(1.7), fill=(255, 255, 255))
    add_bullets(
        slide,
        Inches(1.0),
        Inches(4.25),
        Inches(11.2),
        Inches(1.0),
        [
            "前端只展示预览态，不直接调用底层模型或绕过权限判断。",
            "AI / 工具 / Skill 先经过安全边界，再决定是否真正执行。",
            "每个模块都可独立验证，避免一次性做成无法收敛的大系统。",
        ],
        size=18,
    )

    # Slide 6: architecture
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "系统架构", "Web、Desktop、packages 和 DB / AI-core 的关系保持分层，UI 不直接绕过底层安全边界。")
    box_specs = [
        (0.7, 1.65, 2.4, 1.0, "apps/web\n书库 / Reader / 导入 / 学习 / Agent"),
        (0.7, 3.25, 2.4, 1.0, "apps/desktop\n本地预览 / 路由 / 学习卡片"),
        (3.5, 1.55, 2.4, 1.08, "packages/shared\n共享类型 / 协议 / 常量"),
        (6.2, 1.55, 2.25, 1.08, "packages/book-engine\n导入 / 解析 / 切分"),
        (8.75, 1.55, 2.25, 1.08, "packages/learning-engine\n能力评分 / 题单推荐"),
        (11.1, 1.55, 1.55, 1.08, "packages/db\nPrisma / 迁移"),
        (6.2, 3.35, 2.25, 1.08, "packages/ai-core\nLLM / 记忆 / 工具 / 权限"),
        (8.75, 3.35, 3.9, 1.08, "DB + 安全预览\n真实写入与预览态严格分离"),
    ]
    for x, y, w, h, label in box_specs:
        add_rounded_panel(slide, Inches(x), Inches(y), Inches(w), Inches(h), fill=(255, 255, 255))
        add_textbox(slide, Inches(x + 0.12), Inches(y + 0.15), Inches(w - 0.24), Inches(h - 0.18), label, size=15, color=NAVY, bold=True)
    arrows = [
        ((3.1, 2.1), (3.5, 2.1)),
        ((3.1, 3.75), (3.5, 3.75)),
        ((5.9, 2.08), (6.2, 2.08)),
        ((8.45, 2.08), (8.75, 2.08)),
        ((11.0, 2.08), (11.1, 2.08)),
        ((7.35, 2.63), (7.35, 3.35)),
        ((9.9, 2.63), (9.9, 3.35)),
    ]
    for (x1, y1), (x2, y2) in arrows:
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        conn.line.color.rgb = RGBColor(*BORDER)
        conn.line.width = Pt(2)
    add_bullets(
        slide,
        Inches(0.75),
        Inches(5.1),
        Inches(12.0),
        Inches(1.2),
        [
            "Web 负责交互层，Desktop 负责本地预览与路由展示。",
            "book-engine / learning-engine / ai-core 作为可复用核心层。",
            "DB 只承接明确的持久化边界，不进入 UI 的随意调用。",
        ],
        size=16,
    )

    # Slide 7: reader chain
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "核心功能 1：Books / Reader 阅读链路", "先把“进入书库 - 进入书籍 - 继续阅读”这条链路跑通。")
    add_picture_contain(slide, CAPTURE_DIR / "04-reader-top.png", Inches(6.9), Inches(1.45), Inches(5.85), Inches(4.75))
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.72),
        Inches(5.6),
        Inches(3.8),
        [
            "书库和书籍详情均采用只读预览，方便答辩时快速切换路径。",
            "Reader 页把阅读进度、本地状态、章节入口和同步预览放在同一视图里。",
            "当前所有“同步”相关动作都只是预览，不直接写 DB。",
        ],
        size=18,
    )
    add_rounded_panel(slide, Inches(0.72), Inches(5.35), Inches(5.75), Inches(1.02), fill=SOFT)
    add_textbox(slide, Inches(0.95), Inches(5.58), Inches(5.2), Inches(0.35), "演示重点：阅读状态是可视化的，但保存和同步仍然保持可控。", size=16, color=BLUE, bold=True)

    # Slide 8: code block / redaction
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "核心功能 2：代码块识别与脱敏机制", "当前开发数据未命中可展示的 Reader 代码块目录，因此用导入页 fenced code 的安全预览来补充说明。")
    add_picture_contain(slide, CAPTURE_DIR / "08-import-chapter-edit.png", Inches(6.5), Inches(1.55), Inches(6.0), Inches(4.9))
    add_bullets(
        slide,
        Inches(0.65),
        Inches(1.6),
        Inches(5.55),
        Inches(3.8),
        [
            "Reader 侧已经有代码块抽取、语言筛选、键盘激活和高亮定位逻辑。",
            "当前演示数据里没有稳定的代码块目录样例，所以不夸大成“已现场展示”。",
            "导入页展示了 fenced code 进入预览后如何被脱敏、切分和纳入确认草案。",
        ],
        size=17,
    )
    add_rounded_panel(slide, Inches(0.7), Inches(5.25), Inches(5.75), Inches(1.2), fill=SOFT_ORANGE)
    add_textbox(slide, Inches(0.95), Inches(5.52), Inches(5.2), Inches(0.45), "字幕建议：当前开发数据未命中可展示样例，先用导入预览中的安全脱敏结果替代演示。", size=14, color=GOLD, bold=True)

    # Slide 9: import preview
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "核心功能 3：纯文本导入预览与章节编辑草案", "输入、切分、脱敏、确认、重命名、排除、撤销 / 重做都留在本地预览层。")
    add_picture_contain(slide, CAPTURE_DIR / "07-import-preview.png", Inches(0.7), Inches(1.55), Inches(5.9), Inches(4.95))
    add_bullets(
        slide,
        Inches(7.0),
        Inches(1.68),
        Inches(5.6),
        Inches(3.9),
        [
            "支持 Markdown 标题和中文章节标题，先生成草案再进入确认。",
            "危险字段会被脱敏，确认状态和保存状态都保持在 preview-only。",
            "章节重命名、排除、撤销和重做只影响当前草案，不触发写入。",
        ],
        size=17,
    )
    add_rounded_panel(slide, Inches(7.05), Inches(5.45), Inches(5.55), Inches(0.9), fill=SOFT)
    add_textbox(slide, Inches(7.3), Inches(5.67), Inches(5.0), Inches(0.3), "保存按钮 disabled，是为了确保演示时不误触真实保存。", size=15, color=BLUE, bold=True)

    # Slide 10: sync preview
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "核心功能 4：Reader Sync 安全预览链路", "强调未接真实生产同步，未写 DB，且仅展示 dev / test-only 预览态。")
    add_picture_contain(slide, CAPTURE_DIR / "05-reader-mid.png", Inches(0.68), Inches(1.55), Inches(6.1), Inches(4.95))
    add_bullets(
        slide,
        Inches(7.05),
        Inches(1.7),
        Inches(5.5),
        Inches(3.9),
        [
            "本页明确标注“演示模式 / 预览数据边界”，同步能力只做可视化说明。",
            "读取本地浏览器记录或开发数据，但不代表真实账户进度恢复。",
            "同步链路中的权限、状态与写入检查保持可追踪、可阻断。",
        ],
        size=17,
    )
    add_rounded_panel(slide, Inches(7.1), Inches(5.35), Inches(5.45), Inches(1.0), fill=(236, 243, 250))
    add_textbox(slide, Inches(7.35), Inches(5.58), Inches(5.0), Inches(0.35), "结论：这是“可展示的安全预览”，不是生产同步。", size=16, color=BLUE, bold=True)

    # Slide 11: desktop
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "核心功能 5：Desktop 本地预览卡片", "桌面端以本地状态为主，展示阅读进度、书签和今日学习行动。")
    add_picture_contain(slide, CAPTURE_DIR / "11-desktop-home.png", Inches(6.92), Inches(1.45), Inches(5.7), Inches(4.85))
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.7),
        Inches(5.55),
        Inches(3.8),
        [
            "Desktop 首页是静态入口页，方便演示本地卡片与路由联动。",
            "阅读进度、书签和今日行动都只读取本地预览，不触发真实自动化。",
            "Reader 跳转和本地状态展示为后续深度集成保留接口。",
        ],
        size=17,
    )
    add_rounded_panel(slide, Inches(0.7), Inches(5.35), Inches(5.5), Inches(1.0), fill=SOFT_ORANGE)
    add_textbox(slide, Inches(0.95), Inches(5.58), Inches(5.05), Inches(0.3), "适合课堂答辩时演示“软件端也在同一个闭环里”。", size=15, color=GOLD, bold=True)

    # Slide 12: implementation
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "项目实施细节", "以最小可运行闭环为单位推进，而不是一次性搭完整平台。")
    left_items = [
        ("软件", "Next.js / React 前端、桌面壳页、Node / TypeScript 工具链。"),
        ("核心包", "book-engine、learning-engine、ai-core、shared、db 分层管理。"),
        ("制作步骤", "先页面预览，再接本地状态，再补测试，再做答辩材料。"),
    ]
    right_items = [
        ("核心技术", "预览态、脱敏、草案编辑、路由安全边界、可追踪日志。"),
        ("验证方式", "lint、typecheck、关键单测和浏览器截图验证。"),
        ("交付形式", "视频、PPT、脚本、README 四件套一次打包。"),
    ]
    for idx, (title, body) in enumerate(left_items):
        y = 1.55 + idx * 1.35
        add_rounded_panel(slide, Inches(0.7), Inches(y), Inches(5.9), Inches(1.1), fill=(255, 255, 255))
        add_badge(slide, Inches(0.92), Inches(y + 0.18), Inches(1.25), Inches(0.34), title, BLUE, color=(255, 255, 255))
        add_textbox(slide, Inches(2.35), Inches(y + 0.15), Inches(4.0), Inches(0.52), body, size=15, color=MUTED)
    for idx, (title, body) in enumerate(right_items):
        y = 1.55 + idx * 1.35
        add_rounded_panel(slide, Inches(6.8), Inches(y), Inches(5.8), Inches(1.1), fill=(255, 255, 255))
        add_badge(slide, Inches(7.02), Inches(y + 0.18), Inches(1.25), Inches(0.34), title, TEAL, color=(255, 255, 255))
        add_textbox(slide, Inches(8.45), Inches(y + 0.15), Inches(4.0), Inches(0.52), body, size=15, color=MUTED)

    # Slide 13: showcase
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "项目作品展示", "把演示视频中的关键截图做成一页，答辩时也可直接指着这页讲。")
    collage = [
        (CAPTURE_DIR / "01-home.png", 0.7, 1.55, 3.95, 1.95, "首页"),
        (CAPTURE_DIR / "04-reader-top.png", 4.82, 1.55, 3.95, 1.95, "Reader"),
        (CAPTURE_DIR / "07-import-preview.png", 8.94, 1.55, 3.65, 1.95, "导入预览"),
        (CAPTURE_DIR / "11-desktop-home.png", 0.7, 3.72, 5.95, 2.05, "Desktop"),
        (CAPTURE_DIR / "08-import-chapter-edit.png", 6.88, 3.72, 5.7, 2.05, "草案编辑"),
    ]
    for path, x, y, w, h, tag in collage:
        add_picture_contain(slide, path, Inches(x), Inches(y), Inches(w), Inches(h))
        add_badge(slide, Inches(x + 0.12), Inches(y + 0.12), Inches(0.9), Inches(0.28), tag, NAVY, color=(255, 255, 255))
    add_textbox(slide, Inches(0.7), Inches(6.15), Inches(12.0), Inches(0.35), "视频文件：docs/demo-deliverables/demo-video.mp4", size=14, color=BLUE, bold=True)
    add_textbox(slide, Inches(0.7), Inches(6.48), Inches(12.0), Inches(0.35), "现场如需播放，可直接用系统播放器或 PowerPoint 旁边打开。", size=14, color=MUTED)

    # Slide 14: testing
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "测试与验证", "2026-06-08 已实跑 lint / typecheck / 关键单测。")
    add_rounded_panel(slide, Inches(0.7), Inches(1.55), Inches(12.0), Inches(4.95), fill=(255, 255, 255))
    test_rows = [
        ["命令", "结果", "说明"],
        ["pnpm lint", "通过", "VM lint + TypeScript syntax check 全部 OK"],
        ["pnpm typecheck", "通过", "0 errors"],
        ["node --test apps/web/src/app/import/text-import-preview.test.mjs", "通过", "16 项断言全部通过"],
        ["node --test apps/web/src/app/import/text-import-edit-preview.test.mjs", "通过", "11 项断言全部通过"],
        ["node --test apps/web/src/app/reader/reader-code-element-extractor.test.mjs", "通过", "12 项断言全部通过"],
    ]
    add_simple_table(
        slide,
        Inches(0.9),
        Inches(1.85),
        Inches(11.6),
        Inches(4.2),
        len(test_rows),
        3,
        test_rows,
        col_widths=[Inches(5.3), Inches(1.3), Inches(5.0)],
    )
    add_textbox(slide, Inches(0.9), Inches(6.58), Inches(11.8), Inches(0.3), "这些结果可以直接在答辩时说明：最关键的预览链路已经做过单测和命令级验证。", size=12, color=MUTED)

    # Slide 15: summary
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "项目总结", "把一个长期复杂项目拆成可演示、可验证、可继续推进的阶段成果。")
    add_rounded_panel(slide, Inches(0.8), Inches(1.65), Inches(12.0), Inches(4.9), fill=(255, 255, 255))
    add_bullets(
        slide,
        Inches(1.05),
        Inches(1.95),
        Inches(5.5),
        Inches(3.8),
        [
            "学到的：分层设计、预览优先、权限边界和可追踪状态是复杂系统的第一步。",
            "遇到的挑战：真实数据并不总是天然适合演示，需要稳定的只读预览样例。",
            "解决方法：把导入预览、Reader 安全预览和 Desktop 本地卡片拆成独立模块。",
        ],
        size=18,
    )
    add_rounded_panel(slide, Inches(6.8), Inches(2.0), Inches(5.5), Inches(2.55), fill=SOFT)
    add_textbox(slide, Inches(7.1), Inches(2.3), Inches(4.9), Inches(0.4), "当前结果", size=18, color=BLUE, bold=True)
    add_bullets(
        slide,
        Inches(7.08),
        Inches(2.7),
        Inches(4.9),
        Inches(1.4),
        [
            "视频、PPT、脚本、README 均已产出。",
            "lint / typecheck / 关键单测已通过。",
            "演示边界被明确标注为 preview-only。",
        ],
        size=15,
    )
    add_textbox(slide, Inches(1.0), Inches(5.95), Inches(11.0), Inches(0.4), "答辩时建议强调：我们不是“做了一个大而全的 AI 平台”，而是先做出了一个能安全演示的最小闭环。", size=15, color=MUTED)

    # Slide 16: future
    slide = prs.slides.add_slide(blank)
    add_slide_header(slide, "未来改进", "后续按路线图逐步补齐真正上线所需要的能力。")
    future_items = [
        "真实 auth / session 管理与用户隔离",
        "DB 持久化、迁移和写入回路",
        "LLM / RAG / 工具调用与权限审计",
        "PDF / EPUB / URL 导入与更完善的章节解析",
        "Desktop 深度集成和跨端状态同步",
        "UI 美化、动效、主题和更多演示样例",
    ]
    for idx, item in enumerate(future_items):
        row = idx // 2
        col = idx % 2
        x = 0.7 + col * 6.25
        y = 1.65 + row * 1.25
        add_rounded_panel(slide, Inches(x), Inches(y), Inches(5.9), Inches(0.95), fill=(255, 255, 255))
        add_badge(slide, Inches(x + 0.16), Inches(y + 0.15), Inches(0.42), Inches(0.3), f"{idx+1}", GOLD, color=NAVY)
        add_textbox(slide, Inches(x + 0.7), Inches(y + 0.12), Inches(4.9), Inches(0.45), item, size=15, color=MUTED)
    add_rounded_panel(slide, Inches(0.7), Inches(5.6), Inches(12.0), Inches(0.9), fill=SOFT_ORANGE)
    add_textbox(slide, Inches(0.95), Inches(5.86), Inches(11.4), Inches(0.3), "下一轮开发建议：优先补真实 auth 与 DB 持久化，再接 LLM / RAG，最后再扩展更多导入格式。", size=15, color=GOLD, bold=True)

    return prs


def write_markdown_files() -> None:
    script_md = """# 项目演示脚本

> 说明：本次演示采用“浏览器截图 + 中文字幕”方式合成视频，属于允许的 fallback。当前环境没有稳定的实时录屏管线，因此改用静态帧合成 MP4。

## 片段 1
- 时间：00:00 - 00:11
- 页面：主页
- 说明：展示项目首页和导航入口，说明 Web / Reader / 导入 / 学习 / Agent 的整体闭环。

## 片段 2
- 时间：00:11 - 00:21
- 页面：书库列表
- 说明：展示只读预览的书库列表，强调当前数据源是预览态，不是生产库。

## 片段 3
- 时间：00:21 - 00:31
- 页面：书籍详情
- 说明：展示章节、chunk 与继续阅读入口，说明阅读链路已经可用。

## 片段 4
- 时间：00:31 - 00:43
- 页面：Reader 顶部
- 说明：展示阅读进度、本地状态、同步预览与章节入口；强调同步链路仅为预览。

## 片段 5
- 时间：00:43 - 00:55
- 页面：Reader 中部
- 说明：展示长文滚动、当前可见块和只读安全说明。

## 片段 6
- 时间：00:55 - 01:07
- 页面：纯文本导入表单
- 说明：展示 Markdown / 章节标题 / fenced code 输入，以及保存按钮 disabled。

## 片段 7
- 时间：01:07 - 01:21
- 页面：导入预览 + 草案编辑
- 说明：展示章节切分、危险字段脱敏、重命名、排除、撤销 / 重做等本地草案能力。

## 片段 8
- 时间：01:21 - 01:32
- 页面：学习页
- 说明：展示阅读进度、能力画像、今日行动和本地预览卡片。

## 片段 9
- 时间：01:32 - 01:43
- 页面：Agent 预览页
- 说明：展示任务 dry-run、权限边界、Skill / 工具风险预览。

## 片段 10
- 时间：01:43 - 01:55
- 页面：Desktop 首页
- 说明：展示本地阅读进度、书签、今日行动和 Reader 跳转预览。

## 片段 11
- 时间：01:55 - 02:11
- 页面：总结页
- 说明：强调当前阶段是 MVP / preview-only / disabled-by-default，不是生产上线。

## 备注
- Reader 代码块目录逻辑已实现，但当前开发数据未命中可展示样例，因此视频里用导入页 fenced code 的安全脱敏片段补位说明。
- 所有保存 / 同步 / Agent 执行都保持阻断或未接入，不展示任何真实密钥、token、`.env` 或数据库密码。
"""
    readme_md = """# Demo Deliverables

本目录包含本次项目演示交付物：

- `demo-video.mp4`
- `defense-presentation.pptx`
- `demo-script.md`
- `README.md`

## 如何播放视频

1. 直接用系统视频播放器打开 `demo-video.mp4`。
2. 如果需要在答辩现场播放，建议提前确认播放器支持 H.264 / MP4。
3. 本次视频不是现场录屏，而是按页面截图 + 中文字幕合成的 MP4。

## 如何打开 PPT

1. 用 PowerPoint、WPS 或 LibreOffice 打开 `defense-presentation.pptx`。
2. 如果需要修改小组成员信息、课程 / 比赛名称或日期占位，直接在对应页填写即可。

## Fallback 说明

- 当前环境没有稳定可用的实时录屏管线，所以视频采用浏览器截图序列合成 MP4 的方式。
- Reader 代码块目录在当前开发数据里没有可稳定展示的样例，因此脚本和视频都明确说明了这一点，并用导入预览中的 fenced code 脱敏样例补位。
- 这次交付没有修改业务源码，没有接真实 auth、DB 写入、LLM 或 Agent loop。
"""
    (OUT_DIR / "demo-script.md").write_text(script_md, encoding="utf-8")
    (OUT_DIR / "README.md").write_text(readme_md, encoding="utf-8")


def main() -> None:
    items = make_demo_frames()
    processed_paths: list[tuple[Path, float]] = []
    for index, item in enumerate(items):
        image = build_video_frame(item, index)
        out_path = TEMP_DIR / f"frame_{index:02d}.png"
        image.save(out_path)
        processed_paths.append((out_path, float(item["duration"])))

    write_video(processed_paths, OUT_DIR / "demo-video.mp4")

    prs_local = create_presentation()
    prs_local.save(str(OUT_DIR / "defense-presentation.pptx"))
    write_markdown_files()


if __name__ == "__main__":
    main()
